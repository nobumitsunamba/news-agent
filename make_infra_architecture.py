"""
生成AI日報エージェント — IT基盤アーキテクチャ図（サービス・インフラ全体）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import pptx.oxml.ns as nsmap
from lxml import etree

# ── カラー定義 ────────────────────────────────
def rgb(r,g,b): return RGBColor(r,g,b)

C_BG        = rgb(0x0E,0x17,0x24)
C_PANEL     = rgb(0x14,0x22,0x33)
C_WHITE     = rgb(0xFF,0xFF,0xFF)
C_MUTED     = rgb(0x7A,0x98,0xB5)
C_CODE      = rgb(0xF8,0xD2,0x66)
C_BORDER    = rgb(0x25,0x3F,0x58)

# プラットフォームカラー
C_GITHUB    = rgb(0x6E,0x40,0xC9)   # GitHub 紫
C_GH_LIGHT  = rgb(0x2A,0x18,0x50)   # GitHub パネル背景
C_ANTHROPIC = rgb(0xCC,0x4B,0x00)   # Anthropic オレンジ
C_ANT_LIGHT = rgb(0x38,0x18,0x00)
C_GOOGLE    = rgb(0x34,0xA8,0x53)   # Google 緑
C_GOO_LIGHT = rgb(0x0A,0x28,0x16)
C_ITINFRA   = rgb(0x00,0x8B,0xD4)   # RSS/外部 青
C_ITA_LIGHT = rgb(0x00,0x18,0x35)
C_RUNNER    = rgb(0x00,0xC8,0x8A)   # Runner / Python 緑
C_RUN_LIGHT = rgb(0x00,0x20,0x14)
C_USER      = rgb(0xFF,0xCC,0x00)   # エンドユーザー

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = C_BG
sh = slide.shapes

# ── ユーティリティ ────────────────────────────
def box(l,t,w,h, fill, lc=None, lw=Pt(0.8), r=None):
    s = sh.add_shape(1, Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if lc: s.line.color.rgb=lc; s.line.width=lw
    else:  s.line.fill.background()
    if r:
        pg = s._element.find('.//' + nsmap.qn('a:prstGeom'))
        if pg is not None:
            pg.set('prst','roundRect')
            av = pg.find(nsmap.qn('a:avLst'))
            if av is None: av=etree.SubElement(pg,nsmap.qn('a:avLst'))
            gd=etree.SubElement(av,nsmap.qn('a:gd'))
            gd.set('name','adj'); gd.set('fmla',f'val {r}')
    return s

def txt(shape, lines, sizes, bolds, colors, align=PP_ALIGN.LEFT, va=None):
    tf = shape.text_frame; tf.word_wrap=True
    if va: tf.vertical_anchor=va
    for i,(t,s,b,c) in enumerate(zip(lines,sizes,bolds,colors)):
        p = tf.paragraphs[i] if i==0 else tf.add_paragraph()
        p.alignment=align
        r=p.add_run(); r.text=t; r.font.size=Pt(s); r.font.bold=b
        r.font.color.rgb=c
        r.font.name='Consolas' if c==C_CODE else 'Meiryo UI'

def tb(l,t,w,h, lines,sizes,bolds,colors, align=PP_ALIGN.LEFT):
    bx=sh.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    txt(bx,lines,sizes,bolds,colors,align); return bx

def lbl(l,t,w,h, text, sz=7.5, c=C_MUTED, align=PP_ALIGN.CENTER, code=False):
    bx=sh.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame; p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.color.rgb=c
    r.font.name='Consolas' if code else 'Meiryo UI'

def arr(x1,y1,x2,y2, c=C_MUTED, w=Pt(1.5), both=False):
    cn=sh.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=c; cn.line.width=w
    ln=cn.line._ln
    for tag,side in [('headEnd','head'),('tailEnd','tail')]:
        e=ln.find(nsmap.qn(f'a:{tag}'))
        if e is None: e=etree.SubElement(ln,nsmap.qn(f'a:{tag}'))
        if side=='head' or (side=='tail' and both):
            e.set('type','arrow'); e.set('w','med'); e.set('len','med')
        else:
            e.set('type','none')
    return cn

def zone_header(l,t,w, text, c):
    """ゾームタイトルバー"""
    s=box(l,t,w,0.28, fill=c, r=8000)
    txt(s,[text],[8],[True],[C_WHITE],align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# タイトルバー
# ══════════════════════════════════════════════
s=box(0,0,13.33,0.48, fill=rgb(0x08,0x10,0x1A))
txt(s,
    ["生成AI日報エージェント  —  IT基盤アーキテクチャ（サービス・インフラ全体）"],
    [17],[True],[C_WHITE], align=PP_ALIGN.CENTER, va=MSO_ANCHOR.MIDDLE)
lbl(10.5,0.06,2.7,0.35,"As of 2026-05-31",sz=8,c=C_MUTED,align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════
# 5つのプラットフォームゾーン
# ══════════════════════════════════════════════
# ── ゾーン① 外部データソース（左端）
box(0.08,0.56,2.55,6.72, fill=C_ITA_LIGHT, lc=C_ITINFRA, lw=Pt(1.0), r=10000)
zone_header(0.08,0.56,2.55,"External Data Sources",C_ITINFRA)

# ── ゾーン② GitHub Platform（左中）
box(2.75,0.56,3.55,6.72, fill=C_GH_LIGHT, lc=C_GITHUB, lw=Pt(1.2), r=10000)
zone_header(2.75,0.56,3.55,"GitHub Platform",C_GITHUB)

# ── ゾーン③ Execution Runtime（中）
box(6.42,0.56,2.75,6.72, fill=C_RUN_LIGHT, lc=C_RUNNER, lw=Pt(1.2), r=10000)
zone_header(6.42,0.56,2.75,"Execution Runtime",C_RUNNER)

# ── ゾーン④ Anthropic Platform（右中）
box(9.28,0.56,2.05,3.45, fill=C_ANT_LIGHT, lc=C_ANTHROPIC, lw=Pt(1.0), r=10000)
zone_header(9.28,0.56,2.05,"Anthropic",C_ANTHROPIC)

# ── ゾーン⑤ Google Platform（右）
box(11.42,0.56,1.83,6.72, fill=C_GOO_LIGHT, lc=C_GOOGLE, lw=Pt(1.0), r=10000)
zone_header(11.42,0.56,1.83,"Google",C_GOOGLE)

# ── ゾーン⑥ Anthropic下（End User）
box(9.28,4.1,2.05,3.18, fill=rgb(0x18,0x18,0x10), lc=C_USER, lw=Pt(1.0), r=10000)
zone_header(9.28,4.1,2.05,"End User",C_USER)

# ══════════════════════════════════════════════
# ゾーン① 外部データソース
# ══════════════════════════════════════════════
rss_sources = [
    ("Google News",   "RSS 2.0 / HTTPS",  "q=生成AI (UTF-8 encoded)",    1.08),
    ("Google News",   "RSS 2.0 / HTTPS",  "q=ChatGPT OR Claude OR...",   2.05),
    ("ITmedia NEWS",  "RSS 2.0 / HTTPS",  "rss.itmedia.co.jp",           3.02),
    ("TechCrunch JP", "RSS 2.0 / HTTPS",  "jp.techcrunch.com",           3.99),
]
for name, proto, url, top in rss_sources:
    s=box(0.16,top,2.38,0.75, fill=rgb(0x00,0x22,0x42), lc=C_ITINFRA, lw=Pt(0.6), r=8000)
    txt(s,[name,proto,url],[9,7,7],[True,False,False],
        [C_ITINFRA,C_MUTED,C_CODE])

lbl(0.16,5.1,2.38,0.28,"Format: XML / Atom/RSS",sz=7,c=C_MUTED,code=True)
lbl(0.16,5.38,2.38,0.28,"Parse: feedparser 6.x",sz=7,c=C_CODE,code=True)
lbl(0.16,5.65,2.38,0.55,"Collect: 前日 00:00〜24:00 JST\nFilter: 17 AI keywords",sz=7,c=C_MUTED)

# ══════════════════════════════════════════════
# ゾーン② GitHub Platform
# ══════════════════════════════════════════════
# Repository
s=box(2.85,1.0,3.35,0.85, fill=rgb(0x1A,0x0A,0x38), lc=C_GITHUB, lw=Pt(0.8), r=8000)
txt(s,["Repository","nobumitsunamba/news-agent","Branch: main  |  Visibility: Public"],
    [9,8,7],[True,False,False],[C_GITHUB,C_MUTED,C_CODE])

# GitHub Actions Scheduler
s=box(2.85,2.02,3.35,1.0, fill=rgb(0x1E,0x0C,0x45), lc=C_GITHUB, lw=Pt(0.8), r=8000)
txt(s,["Actions Scheduler (Cron)",
       "cron: '0 22 * * *'  (UTC)",
       "= JST 07:00 daily",
       "workflow_dispatch (手動)"],
    [9,8,8,7],[True,False,False,False],
    [C_GITHUB,C_CODE,C_MUTED,C_MUTED])

# GitHub Actions Runner
s=box(2.85,3.18,3.35,1.2, fill=rgb(0x1A,0x0A,0x38), lc=C_GITHUB, lw=Pt(0.8), r=8000)
txt(s,["Actions Runner",
       "ubuntu-latest (ephemeral VM)",
       "timeout: 15 min",
       "Python 3.12  |  pip cache",
       "pip install -r requirements.txt"],
    [9,8,7,7,7],[True,False,False,False,False],
    [C_GITHUB,C_MUTED,C_MUTED,C_CODE,C_CODE])

# GitHub Secrets
s=box(2.85,4.55,3.35,1.55, fill=rgb(0x14,0x08,0x30), lc=C_GITHUB, lw=Pt(0.8), r=8000)
txt(s,["Secrets Manager",
       "ANTHROPIC_API_KEY   ••••••",
       "GMAIL_ADDRESS          ••••••",
       "GMAIL_APP_PASSWORD  ••••••",
       "TO_EMAIL                    ••••••",
       "→ env vars に注入して実行"],
    [9,7.5,7.5,7.5,7.5,7.5],[True,False,False,False,False,False],
    [rgb(0xA0,0x80,0xFF),C_CODE,C_CODE,C_CODE,C_CODE,C_MUTED])

# ══════════════════════════════════════════════
# ゾーン③ Execution Runtime
# ══════════════════════════════════════════════
# Python Runtime
s=box(6.52,1.0,2.55,0.72, fill=rgb(0x00,0x28,0x18), lc=C_RUNNER, lw=Pt(0.8), r=8000)
txt(s,["Python 3.12 Runtime",
       "CPython  |  ubuntu-latest",
       "pip + cache (actions/setup-python@v5)"],
    [9,7.5,7],[True,False,False],[C_RUNNER,C_MUTED,C_MUTED])

# news_agent.py
s=box(6.52,1.88,2.55,3.8, fill=rgb(0x00,0x22,0x14), lc=C_RUNNER, lw=Pt(1.0), r=8000)
txt(s,["news_agent.py",
       "",
       "collect_all_articles()",
       "  feedparser.parse(url)",
       "  → list[dict]",
       "",
       "filter_articles()",
       "  keyword match (lower)",
       "  → list[dict]",
       "",
       "summarize_with_claude()",
       "  client.messages.create()",
       "  → str (markdown)",
       "",
       "send_email()",
       "  SMTP_SSL :465",
       "  MIMEMultipart"],
    [9,7,7.5,7,7,7,7.5,7,7,7,7.5,7,7,7,7.5,7,7],
    [True,False,True,False,False,False,True,False,False,False,True,False,False,False,True,False,False],
    [C_RUNNER,C_MUTED,C_CODE,C_MUTED,C_MUTED,C_MUTED,C_CODE,C_MUTED,C_MUTED,C_MUTED,
     C_CODE,C_MUTED,C_MUTED,C_MUTED,C_CODE,C_MUTED,C_MUTED])

# Libraries
s=box(6.52,5.85,2.55,1.1, fill=rgb(0x00,0x1C,0x10), lc=C_BORDER, lw=Pt(0.5), r=8000)
txt(s,["requirements.txt",
       "anthropic >= 0.40.0",
       "feedparser >= 6.0.11",
       "python-dotenv >= 1.0.0"],
    [8,7.5,7.5,7.5],[True,False,False,False],
    [C_MUTED,C_CODE,C_CODE,C_CODE])

# ══════════════════════════════════════════════
# ゾーン④ Anthropic Platform
# ══════════════════════════════════════════════
s=box(9.38,1.0,1.85,1.3, fill=rgb(0x40,0x18,0x00), lc=C_ANTHROPIC, lw=Pt(0.8), r=8000)
txt(s,["Claude API",
       "claude-sonnet-4-6",
       "POST /v1/messages",
       "max_tokens: 4096"],
    [9,8,7.5,7.5],[True,False,False,False],
    [C_ANTHROPIC,C_CODE,C_MUTED,C_MUTED])

s=box(9.38,2.5,1.85,0.85, fill=rgb(0x28,0x10,0x00), lc=C_BORDER, lw=Pt(0.5), r=8000)
txt(s,["Prompt Task",
       "重複排除 / 重要度ソート",
       "3行日本語要約"],
    [8,7.5,7.5],[True,False,False],[C_MUTED,C_MUTED,C_MUTED])

s=box(9.38,3.5,1.85,0.38, fill=rgb(0x20,0x0C,0x00), lc=C_BORDER, lw=Pt(0.5), r=8000)
txt(s,["HTTPS / JSON  (REST API)"],[7.5],[False],[C_CODE])

# ══════════════════════════════════════════════
# ゾーン⑤ Google Platform
# ══════════════════════════════════════════════
# Google News RSS サービス
s=box(11.52,1.0,1.63,1.3, fill=rgb(0x06,0x24,0x12), lc=C_GOOGLE, lw=Pt(0.8), r=8000)
txt(s,["Google News",
       "RSS Service",
       "news.google.com",
       "/rss/search?q=…"],
    [9,8,7.5,7.5],[True,False,False,False],
    [C_GOOGLE,C_MUTED,C_MUTED,C_CODE])

# Gmail SMTP Gateway
s=box(11.52,2.5,1.63,1.5, fill=rgb(0x06,0x2A,0x14), lc=C_GOOGLE, lw=Pt(0.8), r=8000)
txt(s,["Gmail SMTP",
       "smtp.gmail.com",
       "Port: 465 / SSL",
       "Auth: App Password",
       "(16-digit)"],
    [9,8,7.5,7.5,7.5],[True,False,False,False,False],
    [C_GOOGLE,C_CODE,C_CODE,C_MUTED,C_MUTED])

# ══════════════════════════════════════════════
# ゾーン⑥ End User
# ══════════════════════════════════════════════
s=box(9.38,4.55,1.85,0.9, fill=rgb(0x28,0x28,0x00), lc=C_USER, lw=Pt(0.8), r=8000)
txt(s,["受信トレイ",
       "件名: 【生成AI日報】",
       "YYYY年MM月DD日"],
    [9,7.5,7.5],[True,False,False],[C_USER,C_MUTED,C_CODE])

s=box(9.38,5.6,1.85,0.75, fill=rgb(0x20,0x20,0x00), lc=C_BORDER, lw=Pt(0.5), r=8000)
txt(s,["text/plain + text/html",
       "Charset: UTF-8",
       "Encoding: quoted-printable"],
    [7.5,7.5,7.5],[False,False,False],[C_CODE,C_CODE,C_MUTED])

# ══════════════════════════════════════════════
# 矢印・接続線
# ══════════════════════════════════════════════
# RSS → Runtime (HTTPS)
for y in [1.45, 2.42, 3.39, 4.37]:
    arr(2.54,y, 2.85,y, c=C_ITINFRA, w=Pt(1.5))

# Cron → Runner（トリガー）
arr(4.52,3.02, 4.52,3.18, c=C_GITHUB, w=Pt(1.8))

# Runner → news_agent（起動）
arr(5.35,3.78, 6.52,3.5, c=C_RUNNER, w=Pt(1.5))

# Secrets → Runtime（env injection）
arr(6.17,5.32, 6.52,4.5, c=rgb(0xA0,0x80,0xFF), w=Pt(1.5))

# Runtime → Anthropic API（HTTPS POST）
arr(9.07,2.5, 9.38,1.85, c=C_ANTHROPIC, w=Pt(2.0))
arr(9.38,2.1, 9.07,2.8, c=C_ANTHROPIC, w=Pt(1.5))

# Runtime → Gmail SMTP
arr(9.07,4.6, 11.52,3.4, c=C_GOOGLE, w=Pt(2.0))

# Gmail → End User
arr(10.3,5.0, 10.3,4.55, c=C_USER, w=Pt(1.8))

# Google News → RSS Data Sources（双方向表記）
arr(11.52,1.45, 11.25,1.45, c=C_GOOGLE, w=Pt(1.2))

# ── ラベル
lbl(9.0,2.2,0.4,0.3,"REQ",sz=7,c=C_ANTHROPIC,code=True)
lbl(9.0,2.7,0.4,0.3,"RES",sz=7,c=C_ANTHROPIC,code=True)
lbl(9.7,4.2,1.5,0.3,"SMTP SSL :465",sz=7,c=C_GOOGLE,code=True)
lbl(6.0,5.05,0.7,0.28,"env",sz=7,c=rgb(0xA0,0x80,0xFF),code=True)
lbl(4.3,1.55,1.0,0.28,"trigger",sz=7,c=C_GITHUB)
lbl(5.8,3.3,0.7,0.28,"exec",sz=7,c=C_RUNNER)

# RSS feeds のプロトコルラベル
for y in [1.45,2.42,3.39,4.37]:
    lbl(2.52,y-0.22,0.5,0.22,"HTTPS",sz=6.5,c=C_ITINFRA,code=True)

# ══════════════════════════════════════════════
# フッター
# ══════════════════════════════════════════════
box(0,7.18,13.33,0.32, fill=rgb(0x08,0x10,0x1A), lc=C_BORDER, lw=Pt(0.3))
legend = [
    ("■ External Data Sources", C_ITINFRA,  0.2),
    ("■ GitHub Platform",       C_GITHUB,   2.6),
    ("■ Execution Runtime",     C_RUNNER,   4.8),
    ("■ Anthropic Platform",    C_ANTHROPIC, 6.9),
    ("■ Google Platform",       C_GOOGLE,   9.0),
    ("■ End User",              C_USER,    11.0),
]
for t,c,l in legend:
    lbl(l,7.21,2.1,0.28,t,sz=7.5,c=c)

# ══════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════
out = "/home/user/news-agent/生成AI日報エージェント_IT基盤アーキテクチャ.pptx"
prs.save(out)
print(f"✅ 保存完了: {out}")
