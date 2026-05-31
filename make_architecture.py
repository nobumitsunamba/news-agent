"""
生成AI日報エージェント ITアーキテクチャ図（エンジニア向け PowerPoint 1枚）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.oxml.ns as nsmap
from lxml import etree

# ── カラーパレット ──────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # 濃紺背景
C_PANEL     = RGBColor(0x16, 0x28, 0x3D)   # パネル背景
C_BORDER    = RGBColor(0x2A, 0x45, 0x65)   # 枠線
C_BLUE      = RGBColor(0x29, 0x8F, 0xFF)   # RSSフィード
C_GREEN     = RGBColor(0x00, 0xC8, 0x7A)   # 処理ステップ
C_TEAL      = RGBColor(0x00, 0xC8, 0xC8)   # Claude API
C_ORANGE    = RGBColor(0xFF, 0x7A, 0x00)   # Gmail
C_PURPLE    = RGBColor(0xA0, 0x5C, 0xFF)   # GitHub Actions
C_GRAY      = RGBColor(0x64, 0x8C, 0xAA)   # Secrets
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED     = RGBColor(0x8A, 0xA8, 0xC0)   # サブテキスト
C_CODE      = RGBColor(0xF8, 0xD2, 0x66)   # コード・関数名
C_DARK_TEXT = RGBColor(0x0D, 0x1B, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── 背景 ────────────────────────────────────────
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = C_BG

shapes = slide.shapes

# ── ヘルパー ────────────────────────────────────
def rgb_hex(r, g, b):
    return RGBColor(r, g, b)

def rect(l, t, w, h, fill, line_c=None, line_w=Pt(1), radius=None):
    s = shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_c:
        s.line.color.rgb = line_c
        s.line.width = line_w
    else:
        s.line.fill.background()
    if radius:
        sp = s._element
        pg = sp.find('.//' + nsmap.qn('a:prstGeom'))
        if pg is not None:
            pg.set('prst', 'roundRect')
            av = pg.find(nsmap.qn('a:avLst'))
            if av is None:
                av = etree.SubElement(pg, nsmap.qn('a:avLst'))
            gd = etree.SubElement(av, nsmap.qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return s

def text_on(shape, lines, sizes, bolds, colors, align=PP_ALIGN.LEFT, v_anchor=None):
    """shape にマルチラインテキストを設定する"""
    from pptx.enum.text import MSO_ANCHOR
    tf = shape.text_frame
    tf.word_wrap = True
    if v_anchor:
        tf.vertical_anchor = v_anchor
    for i, (txt, sz, bd, col) in enumerate(zip(lines, sizes, bolds, colors)):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = col
        run.font.name = "Consolas" if col == C_CODE else "Meiryo UI"

def tb(l, t, w, h, lines, sizes, bolds, colors, align=PP_ALIGN.LEFT):
    box = shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    text_on(box, lines, sizes, bolds, colors, align)
    return box

def arrow(x1, y1, x2, y2, color=C_MUTED, width=Pt(1.5), dash=None):
    c = shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = width
    if dash:
        c.line.dash_style = dash
    ln = c.line._ln
    he = ln.find(nsmap.qn('a:headEnd'))
    if he is None:
        he = etree.SubElement(ln, nsmap.qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')
    te = ln.find(nsmap.qn('a:tailEnd'))
    if te is None:
        te = etree.SubElement(ln, nsmap.qn('a:tailEnd'))
    te.set('type', 'none')
    return c

def label(l, t, w, h, txt, sz=7.5, color=C_MUTED, align=PP_ALIGN.CENTER, code=False):
    box = shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(sz)
    run.font.color.rgb = color
    run.font.name = "Consolas" if code else "Meiryo UI"

# ══════════════════════════════════════════════
# タイトルバー
# ══════════════════════════════════════════════
title_bar = rect(0, 0, 13.33, 0.52, fill=C_PANEL, line_c=C_BORDER, line_w=Pt(0.5))
tb(0.15, 0.07, 9.0, 0.38,
   ["生成AI日報エージェント  —  ITアーキテクチャ（エンジニア向け）"],
   [17], [True], [C_WHITE])
tb(10.0, 0.07, 3.0, 0.38,
   ["Python 3.12  |  GitHub Actions"],
   [9], [False], [C_MUTED], align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════
# ゾーンパネル
# ══════════════════════════════════════════════
# 左ゾーン
z_left = rect(0.12, 0.62, 2.85, 6.2, fill=C_PANEL,
              line_c=C_BLUE, line_w=Pt(0.8), radius=12000)
tb(0.22, 0.65, 2.5, 0.28,
   ["INPUT  /  RSS Feeds"],
   [7.5], [True], [C_BLUE])

# 中央ゾーン
z_mid = rect(3.2, 0.62, 4.3, 6.2, fill=C_PANEL,
             line_c=C_GREEN, line_w=Pt(0.8), radius=12000)
tb(3.3, 0.65, 3.8, 0.28,
   ["PROCESSING  /  news_agent.py"],
   [7.5], [True], [C_GREEN])

# 右ゾーン
z_right = rect(9.72, 0.62, 3.5, 6.2, fill=C_PANEL,
               line_c=C_ORANGE, line_w=Pt(0.8), radius=12000)
tb(9.82, 0.65, 3.0, 0.28,
   ["OUTPUT  /  Gmail SMTP"],
   [7.5], [True], [C_ORANGE])

# ══════════════════════════════════════════════
# 左：RSSフィード 4個
# ══════════════════════════════════════════════
feeds = [
    ("Google News",    "q=生成AI&hl=ja&gl=JP",           1.05),
    ("Google News",    "q=ChatGPT OR Claude OR Gemini OR LLM", 1.85),
    ("ITmedia NEWS",   "rss.itmedia.co.jp/rss/2.0/…",    2.8),
    ("TechCrunch JP",  "jp.techcrunch.com/feed/",         3.6),
]
for name, url, top in feeds:
    s = rect(0.22, top, 2.65, 0.65, fill=rgb_hex(0x0A,0x35,0x60),
             line_c=C_BLUE, line_w=Pt(0.8), radius=10000)
    text_on(s,
            [name, url],
            [9, 7],
            [True, False],
            [C_BLUE, C_MUTED])

# feedparser ラベル
label(0.22, 4.4, 2.65, 0.28,
      "lib: feedparser 6.x  (RSS/Atom parser)",
      sz=7, color=C_CODE, code=True)

# キーワード
kw_box = rect(0.22, 4.72, 2.65, 1.9, fill=rgb_hex(0x0A,0x28,0x40),
              line_c=C_GRAY, line_w=Pt(0.5), radius=8000)
text_on(kw_box,
        ["Keyword Filter  (17 words)",
         "GPT / 生成AI / LLM / 大規模言語モデル",
         "OpenAI / ChatGPT / Anthropic / Claude",
         "Gemini / Grok / DeepSeek / Copilot",
         "IBM BoB / AIエージェント / AI Agent",
         "Agentic AI / RAG"],
        [7.5, 7, 7, 7, 7, 7],
        [True, False, False, False, False, False],
        [C_MUTED, C_CODE, C_CODE, C_CODE, C_CODE, C_CODE])

# ══════════════════════════════════════════════
# 中央：処理ステップ
# ══════════════════════════════════════════════
steps = [
    ("① collect_all_articles()",
     ["RSS fetch  →  feedparser.parse(url)",
      "期間: 前日 00:00〜24:00 JST",
      "戻り値: list[dict]  {title, url, summary}"],
     1.05),
    ("② filter_articles()",
     ["is_ai_related() でキーワード照合",
      "title + summary を lower() 比較",
      "戻り値: list[dict]  (絞り込み済み)"],
     2.2),
    ("③ summarize_with_claude()",
     ["client.messages.create()",
      "model: claude-sonnet-4-6",
      "max_tokens: 4096",
      "→ 重複排除 / 重要度ソート / 3行要約"],
     3.35),
    ("④ send_email()",
     ["smtplib.SMTP_SSL('smtp.gmail.com', 465)",
      "MIMEMultipart('alternative')",
      "text/plain + text/html (UTF-8)"],
     4.65),
]

for title, details, top in steps:
    h = 0.95 if len(details) == 3 else 1.0
    s = rect(3.3, top, 4.1, h, fill=rgb_hex(0x07,0x30,0x20),
             line_c=C_GREEN, line_w=Pt(0.8), radius=10000)
    text_on(s,
            [title] + details,
            [9] + [7.5] * len(details),
            [True] + [False] * len(details),
            [C_GREEN] + [C_MUTED] * len(details))

# ステップ間の矢印
for y in [2.0, 3.15, 4.3]:
    arrow(5.35, y, 5.35, y + 0.2, color=C_GREEN, width=Pt(1.5))

# ══════════════════════════════════════════════
# 右：Gmail / 出力仕様
# ══════════════════════════════════════════════
# Gmail
gm = rect(9.82, 1.05, 3.3, 0.9, fill=rgb_hex(0x40,0x20,0x00),
          line_c=C_ORANGE, line_w=Pt(1.2), radius=10000)
text_on(gm,
        ["Gmail  SMTP",
         "smtp.gmail.com  :  465  (SSL)",
         "Auth: App Password (16-digit)"],
        [11, 8, 8],
        [True, False, False],
        [C_ORANGE, C_MUTED, C_MUTED])

# メール仕様
spec = rect(9.82, 2.1, 3.3, 1.55, fill=rgb_hex(0x28,0x16,0x00),
            line_c=C_BORDER, line_w=Pt(0.5), radius=8000)
text_on(spec,
        ["Mail Spec",
         "Subject : 【生成AI日報】YYYY年MM月DD日",
         "Body    : text/plain + text/html",
         "Charset : UTF-8",
         "Trigger : cron JST 07:00 daily"],
        [8, 7.5, 7.5, 7.5, 7.5],
        [True, False, False, False, False],
        [C_MUTED, C_CODE, C_CODE, C_CODE, C_CODE])

# 環境変数
env = rect(9.82, 3.82, 3.3, 1.55, fill=rgb_hex(0x12,0x20,0x30),
           line_c=C_GRAY, line_w=Pt(0.8), radius=8000)
text_on(env,
        ["GitHub Secrets  (env injection)",
         "ANTHROPIC_API_KEY",
         "GMAIL_ADDRESS",
         "GMAIL_APP_PASSWORD",
         "TO_EMAIL"],
        [8, 7.5, 7.5, 7.5, 7.5],
        [True, False, False, False, False],
        [C_GRAY, C_CODE, C_CODE, C_CODE, C_CODE])

# ══════════════════════════════════════════════
# Claude API ブロック（中央下）
# ══════════════════════════════════════════════
cl = rect(3.3, 5.82, 4.1, 0.85, fill=rgb_hex(0x00,0x28,0x30),
          line_c=C_TEAL, line_w=Pt(1.0), radius=10000)
text_on(cl,
        ["Anthropic Claude API",
         "POST https://api.anthropic.com/v1/messages",
         "lib: anthropic>=0.40.0  |  model: claude-sonnet-4-6"],
        [10, 7.5, 7.5],
        [True, False, False],
        [C_TEAL, C_MUTED, C_CODE])

# ══════════════════════════════════════════════
# GitHub Actions ブロック（上部・横断）
# ══════════════════════════════════════════════
ga = rect(3.3, 0.62, 4.1, 0.55, fill=rgb_hex(0x1A,0x0A,0x35),
          line_c=C_PURPLE, line_w=Pt(1.0), radius=10000)
text_on(ga,
        ["GitHub Actions  |  ubuntu-latest  |  Python 3.12  |  cron: '0 22 * * *' UTC  =  JST 07:00"],
        [8],
        [True],
        [C_PURPLE])

# ══════════════════════════════════════════════
# 矢印
# ══════════════════════════════════════════════
# RSS → ①収集
for y in [1.37, 2.17, 3.12, 3.92]:
    arrow(2.87, y, 3.3, y, color=C_BLUE, width=Pt(1.5))

# ③要約 ↔ Claude API（双方向）
arrow(5.35, 5.82, 5.35, 4.3, color=C_TEAL, width=Pt(1.5))   # 下→上（レスポンス）
arrow(5.35, 4.3, 5.35, 5.82, color=C_TEAL, width=Pt(1.2))   # 上→下（リクエスト）

# ④送信 → Gmail
arrow(7.4, 5.12, 9.82, 1.8, color=C_ORANGE, width=Pt(2.0))

# GitHub Actions → ①（トリガー）
arrow(5.35, 1.17, 5.35, 1.05, color=C_PURPLE, width=Pt(1.8))

# Secrets → ④（env injection）
arrow(9.82, 4.6, 7.4, 4.95, color=C_GRAY, width=Pt(1.2))

# ── ラベル ─────────────────────────────────────
label(6.5, 5.0, 3.0, 0.28,
      "MIME multipart email", sz=7.5, color=C_ORANGE)
label(7.55, 4.35, 2.0, 0.28,
      "env vars inject", sz=7, color=C_GRAY)
label(5.55, 5.3, 1.8, 0.28,
      "HTTP/1.1 POST\nHTTP/1.1 200 OK", sz=7, color=C_TEAL)
label(5.55, 1.1, 2.2, 0.28,
      "workflow_dispatch / schedule", sz=7, color=C_PURPLE)

# ══════════════════════════════════════════════
# フッター（凡例）
# ══════════════════════════════════════════════
rect(0, 7.18, 13.33, 0.32, fill=rgb_hex(0x0A,0x14,0x20),
     line_c=C_BORDER, line_w=Pt(0.3))

legend = [
    ("■ RSS Input",        C_BLUE,   0.2),
    ("■ Processing",       C_GREEN,  2.2),
    ("■ Claude API",       C_TEAL,   4.2),
    ("■ Gmail Output",     C_ORANGE, 6.2),
    ("■ GitHub Actions",   C_PURPLE, 8.2),
    ("■ GitHub Secrets",   C_GRAY,  10.5),
    ("■ Code / Config",    C_CODE,  12.2),
]
for txt, col, l in legend:
    label(l, 7.2, 2.0, 0.28, txt, sz=7.5, color=col)

# ══════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════
output = "/home/user/news-agent/生成AI日報エージェント_アーキテクチャ.pptx"
prs.save(output)
print(f"✅ 保存完了: {output}")
